#!/usr/bin/env python3
"""
PowerPoint Generator for CRM Reports
Creates PowerPoint presentations from extracted CRM data
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
import os
import glob
import re
from datetime import datetime
from typing import Dict, List, Any


class PowerPointGenerator:
    """Generate PowerPoint reports from CRM data"""

    # Color palette for charts
    COLORS = ['#2caffe', '#544fc5', '#00e272', '#fe6a35', '#6b8abc', '#d568fb', '#ff6384']

    def __init__(self, slide_config: List[Dict[str, str]]):
        """
        Initialize PowerPoint generator

        Args:
            slide_config: List of slide configurations
        """
        self.slide_config = slide_config
        self.prs = Presentation()

    def find_svg_file(self, folder: str, pattern: str) -> str:
        """
        Find SVG file matching pattern

        Args:
            folder: Folder to search
            pattern: File pattern to match

        Returns:
            Path to SVG file or None
        """
        folder_path = os.path.join(folder)

        # Exact match
        exact_path = os.path.join(folder_path, pattern)
        if os.path.exists(exact_path):
            return exact_path

        # Debug prefix match
        debug_pattern = f"debug_{pattern}"
        debug_path = os.path.join(folder_path, debug_pattern)
        if os.path.exists(debug_path):
            return debug_path

        # Partial match
        search_pattern = pattern.replace('.svg_idx0.svg', '*.svg')
        matches = glob.glob(os.path.join(folder_path, f"*{search_pattern}"))
        if matches:
            return matches[0]

        return None

    def filter_data_by_org(self, data: Dict[str, Any], org: str) -> Dict[str, Any]:
        """
        Filter data by organization

        Args:
            data: Data to filter
            org: Organization name (ISIM or Perspectivia)

        Returns:
            Filtered data
        """
        filtered_data = {}

        if org == 'ISIM':
            for centre, values in data.items():
                if any(keyword in centre.upper() for keyword in ['ISIM']):
                    filtered_data[centre] = values

        elif org == 'Perspectivia':
            for centre, values in data.items():
                if any(keyword in centre.upper() for keyword in ['PERSPECTIV', 'EDUCAE']):
                    filtered_data[centre] = values

        # If no centers found for organization, use all data
        if not filtered_data:
            print(f"  [INFO] No {org} centers found, using all data")
            filtered_data = data

        return filtered_data

    def create_slide_with_data(self, slide_title: str, svg_title: str,
                               data: Dict[str, Any], org_filter: str,
                               slide_number: int = None):
        """
        Create slide with specific data

        Args:
            slide_title: Slide title
            svg_title: SVG chart title
            data: Chart data
            org_filter: Organization filter (ISIM/Perspectivia)
            slide_number: Slide number
        """
        # Filter by organization
        filtered_data = self.filter_data_by_org(data, org_filter)

        print(f"  {org_filter} data: {len(filtered_data)} center(s)")

        # Create slide
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)

        # Title
        title_shape = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(9), Inches(0.8)
        )
        title_frame = title_shape.text_frame
        title_frame.text = slide_title

        title_paragraph = title_frame.paragraphs[0]
        title_paragraph.alignment = PP_ALIGN.CENTER
        title_font = title_paragraph.font
        title_font.name = 'Arial'
        title_font.size = Pt(18)
        title_font.bold = True
        title_font.color.rgb = RGBColor(44, 62, 80)

        # Subtitle - Extract dates from SVG title
        date_match = re.search(r'du (\d{2}/\d{2}/\d{4}) au (\d{2}/\d{2}/\d{4})', svg_title)
        if date_match:
            date_range = f"du {date_match.group(1)} au {date_match.group(2)}"
        else:
            date_range = ""

        subtitle_shape = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.0), Inches(9), Inches(0.3)
        )
        subtitle_frame = subtitle_shape.text_frame
        subtitle_frame.text = date_range if date_range else ""

        subtitle_paragraph = subtitle_frame.paragraphs[0]
        subtitle_paragraph.alignment = PP_ALIGN.CENTER
        subtitle_font = subtitle_paragraph.font
        subtitle_font.name = 'Arial'
        subtitle_font.size = Pt(12)
        subtitle_font.color.rgb = RGBColor(127, 140, 141)

        # Chart or "no data" message
        if filtered_data:
            # Filter centers without data (all values = 0)
            centres_with_data = {}
            for centre, metrics_data in filtered_data.items():
                has_data = any(val > 0 for val in metrics_data.values() if isinstance(val, (int, float)))
                if has_data:
                    centres_with_data[centre] = metrics_data

            if not centres_with_data:
                centres_with_data = filtered_data

            chart_data = CategoryChartData()
            centres = list(centres_with_data.keys())
            chart_data.categories = centres

            print(f"  Filtered centers: {len(centres)} center(s) with data (out of {len(filtered_data)} total)")

            # Metrics
            all_metrics = set()
            for centre_data in centres_with_data.values():
                all_metrics.update(centre_data.keys())

            metrics = list(all_metrics)[:7]  # Max 7 series

            for metric in metrics:
                values = []
                for centre in centres:
                    val = centres_with_data[centre].get(metric, 0)
                    values.append(val if val > 0 else None)
                chart_data.add_series(metric, values)

            # Add chart
            x, y, cx, cy = Inches(0.5), Inches(1.5), Inches(9), Inches(5.5)
            chart_frame = slide.shapes.add_chart(
                XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
            )
            chart = chart_frame.chart

            # Customize
            chart.has_title = True
            chart.chart_title.text_frame.text = slide_title
            chart.has_legend = True
            chart.legend.position = XL_LEGEND_POSITION.BOTTOM

            # Colors and conditional labels
            for i, series in enumerate(chart.series):
                if i < len(self.COLORS):
                    color_hex = self.COLORS[i].replace('#', '')
                    try:
                        r, g, b = int(color_hex[0:2], 16), int(color_hex[2:4], 16), int(color_hex[4:6], 16)
                        series.format.fill.solid()
                        series.format.fill.fore_color.rgb = RGBColor(r, g, b)
                    except:
                        pass

                # Disable labels for specific slides
                slides_without_labels = [3, 4, 5, 6, 8, 9]
                if slide_number in slides_without_labels:
                    series.has_data_labels = False
                else:
                    series.has_data_labels = True
                    data_labels = series.data_labels
                    data_labels.show_value = True
                    data_labels.font.name = 'Arial'
                    data_labels.font.size = Pt(9)

            # Formatting adjustments
            category_axis = chart.category_axis
            category_labels = category_axis.tick_labels
            category_labels.font.name = 'Arial'
            category_labels.font.size = Pt(9)
            category_labels.offset = 100

            cat_tx_pr = category_axis._element.get_or_add_txPr()
            body_pr = cat_tx_pr.bodyPr
            body_pr.set('rot', '-2700000')  # -45 deg
            body_pr.set('vertOverflow', 'ellipsis')
            body_pr.set('wrap', 'square')
            body_pr.set('anchor', 'ctr')

            value_axis = chart.value_axis
            value_labels = value_axis.tick_labels
            value_labels.font.name = 'Arial'
            value_labels.font.size = Pt(9)

            legend = chart.legend
            if legend is not None:
                legend.font.name = 'Arial'
                legend.font.size = Pt(10)
                legend.include_in_layout = False

            print(f"  [OK] Slide '{slide_title}' created with {len(centres)} center(s) and {len(metrics)} metric(s)")

        else:
            # No data for this organization
            no_data_shape = slide.shapes.add_textbox(
                Inches(2), Inches(3), Inches(6), Inches(2)
            )
            no_data_frame = no_data_shape.text_frame
            no_data_frame.text = "No data available for this organization"

            no_data_paragraph = no_data_frame.paragraphs[0]
            no_data_paragraph.alignment = PP_ALIGN.CENTER
            no_data_font = no_data_paragraph.font
            no_data_font.name = 'Arial'
            no_data_font.size = Pt(20)
            no_data_font.color.rgb = RGBColor(190, 100, 100)

            print(f"  [INFO] Slide '{slide_title}' created without data")

        return slide

    def create_title_slide(self, title: str, subtitle: str):
        """Create title slide"""
        title_slide_layout = self.prs.slide_layouts[0]
        title_slide = self.prs.slides.add_slide(title_slide_layout)
        title_obj = title_slide.shapes.title
        subtitle_obj = title_slide.placeholders[1]

        title_obj.text = title
        subtitle_obj.text = subtitle

    def save(self, output_path: str):
        """Save presentation to file"""
        self.prs.save(output_path)
        print(f"\n{'='*70}")
        print(f"[OK] PowerPoint created: {output_path}")
        print(f"[INFO] Configured slides: {len(self.slide_config)}")
        print(f"[INFO] Total slides: {len(self.prs.slides)} (including title)")
        print(f"[INFO] DATA EXTRACTED FROM CRM!")
        print("="*70)


if __name__ == "__main__":
    print("PowerPoint Generator module loaded successfully")
