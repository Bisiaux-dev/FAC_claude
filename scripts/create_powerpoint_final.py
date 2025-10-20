#!/usr/bin/env python3
"""
Générateur PowerPoint Final - 15 slides spécifiques selon configuration
Solution finale avec extraction des VRAIES données des fichiers SVG
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
import os
import glob

# Importer la fonction de parsing
from parse_svg_data import extract_highcharts_data

# Configuration des 12 slides spécifiques (suppression des slides 2, 7, 11)
SLIDE_CONFIG = [
    {
        'file_pattern': 'jourj_graph2.svg',
        'folder': 'downloads_cf',
        'org': 'ISIM',
        'title': 'Engagement collective quotidienne'
    },
    {
        'file_pattern': 'jourj_graph4.svg',
        'folder': 'downloads_cf',
        'org': 'ISIM',
        'title': 'Productivité individuelle quotidienne'
    },
    {
        'file_pattern': 'semaine_precedente_graph3.svg',
        'folder': 'downloads_cf',
        'org': 'ISIM',
        'title': 'Qualité individuelle quotidienne'
    },
    {
        'file_pattern': 'semaine_precedente_graph4.svg',
        'folder': 'downloads_cf',
        'org': 'ISIM',
        'title': 'Productivité collectif hebdomadaire'
    },
    {
        'file_pattern': 'j90_semaine_graph4.svg',
        'folder': 'downloads_cf',
        'org': 'ISIM',
        'title': 'Productivité collectif trimestre'
    },
    {
        'file_pattern': 'cip_j7_graph0.svg',
        'folder': 'downloads_cip',
        'org': 'Perspectivia',
        'title': 'Productivité collectif hebdomadaire'
    },
    {
        'file_pattern': 'cip_j7_graph2.svg',
        'folder': 'downloads_cip',
        'org': 'Perspectivia',
        'title': 'Engagement collectif hebdomadaire'
    },
    {
        'file_pattern': 'cip_j7_graph4.svg',
        'folder': 'downloads_cip',
        'org': 'ISIM',
        'title': 'Productivité collectif hebdomadaire'
    },
    {
        'file_pattern': 'cip_j90_graph0.svg',
        'folder': 'downloads_cip',
        'org': 'ISIM',
        'title': 'Engagement collective quotidienne'
    },
    {
        'file_pattern': 'cip_j90_graph0.svg',
        'folder': 'downloads_cip',
        'org': 'Perspectivia',
        'title': 'Productivité individuelle hebdomadaire'
    },
    {
        'file_pattern': 'cip_jourj_graph0.svg',
        'folder': 'downloads_cip',
        'org': 'Perspectivia',
        'title': 'Productivité collectif quotidienne'
    },
    {
        'file_pattern': 'cip_jourj_graph4.svg',
        'folder': 'downloads_cip',
        'org': 'Perspectivia',
        'title': 'Productivité individuelle quotidienne'
    }
]

def find_svg_file(folder, pattern):
    """Trouve le fichier SVG correspondant au pattern"""
    base_dir = os.path.dirname(os.path.abspath(__file__))
    folder_path = os.path.join(base_dir, '..', folder)
    
    # Recherche exacte
    exact_path = os.path.join(folder_path, pattern)
    if os.path.exists(exact_path):
        return exact_path
    
    # Recherche avec préfixe debug_
    debug_pattern = f"debug_{pattern}"
    debug_path = os.path.join(folder_path, debug_pattern)
    if os.path.exists(debug_path):
        return debug_path
        
    # Recherche partielle
    search_pattern = pattern.replace('.svg_idx0.svg', '*.svg')
    matches = glob.glob(os.path.join(folder_path, f"*{search_pattern}"))
    if matches:
        return matches[0]
    
    return None

def filter_data_by_org(data, org):
    """Filtre les données selon l'organisation"""
    filtered_data = {}
    
    if org == 'ISIM':
        for centre, values in data.items():
            if any(keyword in centre.upper() for keyword in ['ISIM']):
                filtered_data[centre] = values
    
    elif org == 'Perspectivia':
        for centre, values in data.items():
            if any(keyword in centre.upper() for keyword in ['PERSPECTIV', 'EDUCAE']):
                filtered_data[centre] = values
    
    # Si aucun centre trouvé pour l'organisation, ignorer le filtrage
    if not filtered_data:
        print(f"  [INFO] Aucun centre {org} trouvé, utilisation de toutes les données")
        filtered_data = data
    
    return filtered_data

def create_slide_with_specific_data(prs, slide_title, svg_title, data, org_filter, slide_number=None):
    """Créer slide avec données spécifiques selon l'organisation"""
    
    # Filtrer selon ISIM vs Perspectivia
    filtered_data = filter_data_by_org(data, org_filter)
    
    print(f"  Données {org_filter}: {len(filtered_data)} centre(s)")
    
    # Créer slide
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Titre
    title_shape = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(9), Inches(0.8)
    )
    title_frame = title_shape.text_frame
    title_frame.text = slide_title
    
    title_paragraph = title_frame.paragraphs[0]
    title_paragraph.alignment = PP_ALIGN.CENTER
    title_font = title_paragraph.font
    title_font.name = 'Arial'
    title_font.size = Pt(18)
    title_font.bold = True
    title_font.color.rgb = RGBColor(44, 62, 80)
    
    # Sous-titre - Extraire uniquement les dates du titre SVG
    import re
    # Chercher les dates au format "du XX/XX/XXXX au XX/XX/XXXX"
    date_match = re.search(r'du (\d{2}/\d{2}/\d{4}) au (\d{2}/\d{2}/\d{4})', svg_title)
    if date_match:
        date_range = f"du {date_match.group(1)} au {date_match.group(2)}"
    else:
        # Essayer le format avec tiret "- du XX/XX/XXXX au XX/XX/XXXX"
        date_match2 = re.search(r'-\s*du\s+(\d{2}/\d{2}/\d{4})\s+au\s+(\d{2}/\d{2}/\d{4})', svg_title)
        if date_match2:
            date_range = f"du {date_match2.group(1)} au {date_match2.group(2)}"
        else:
            # Si pas de date trouvée, afficher le titre complet sans le préfixe "KPI"
            date_range = re.sub(r'^KPI\s+[^-]+-\s*', '', svg_title) if svg_title else ""

    subtitle_shape = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.0), Inches(9), Inches(0.3)
    )
    subtitle_frame = subtitle_shape.text_frame
    subtitle_frame.text = date_range if date_range else ""

    subtitle_paragraph = subtitle_frame.paragraphs[0]
    subtitle_paragraph.alignment = PP_ALIGN.CENTER
    subtitle_font = subtitle_paragraph.font
    subtitle_font.name = 'Arial'
    subtitle_font.size = Pt(12)
    subtitle_font.color.rgb = RGBColor(127, 140, 141)
    
    # Graphique ou message "pas de données"
    if filtered_data:
        # FILTRER LES CENTRES SANS AUCUNE DONNÉE (toutes valeurs à 0 ou vides)
        centres_with_data = {}
        for centre, metrics_data in filtered_data.items():
            # Vérifier si au moins une métrique a une valeur > 0
            has_data = any(val > 0 for val in metrics_data.values() if isinstance(val, (int, float)))
            if has_data:
                centres_with_data[centre] = metrics_data

        # Si aucun centre n'a de données après filtrage
        if not centres_with_data:
            centres_with_data = filtered_data  # Garder tout pour éviter un graphique vide

        chart_data = CategoryChartData()
        centres = list(centres_with_data.keys())
        chart_data.categories = centres

        print(f"  Centres filtrés: {len(centres)} centre(s) avec données (sur {len(filtered_data)} total)")

        # Métriques
        all_metrics = set()
        for centre_data in centres_with_data.values():
            all_metrics.update(centre_data.keys())

        metrics = list(all_metrics)[:7]  # Max 7 séries
        colors = ['#2caffe', '#544fc5', '#00e272', '#fe6a35', '#6b8abc', '#d568fb', '#ff6384']

        for metric in metrics:
            values = []
            for centre in centres:
                val = centres_with_data[centre].get(metric, 0)
                values.append(val if val > 0 else None)  # None n'affiche pas d'étiquette
            chart_data.add_series(metric, values)
        
        # Ajouter graphique
        x, y, cx, cy = Inches(0.5), Inches(1.5), Inches(9), Inches(5.5)
        chart_frame = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
        )
        chart = chart_frame.chart
        
        # Personnaliser
        chart.has_title = True
        chart.chart_title.text_frame.text = slide_title
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        
        # Couleurs et étiquettes conditionnelles
        for i, series in enumerate(chart.series):
            if i < len(colors):
                color_hex = colors[i].replace('#', '')
                try:
                    r, g, b = int(color_hex[0:2], 16), int(color_hex[2:4], 16), int(color_hex[4:6], 16)
                    series.format.fill.solid()
                    series.format.fill.fore_color.rgb = RGBColor(r, g, b)
                except:
                    pass
            
            # Désactiver les étiquettes pour les slides spécifiées
            slides_without_labels = [3, 4, 5, 6, 8, 9]
            if slide_number in slides_without_labels:
                series.has_data_labels = False
            else:
                # Activer les étiquettes de données (les zéros sont maintenant None donc invisibles)
                series.has_data_labels = True
                data_labels = series.data_labels
                data_labels.show_value = True
                data_labels.font.name = 'Arial'
                data_labels.font.size = Pt(9)

        # Ajustements de mise en forme pour PowerPoint
        category_axis = chart.category_axis
        category_labels = category_axis.tick_labels
        category_labels.font.name = 'Arial'
        category_labels.font.size = Pt(9)
        category_labels.offset = 100

        cat_tx_pr = category_axis._element.get_or_add_txPr()
        body_pr = cat_tx_pr.bodyPr
        body_pr.set('rot', '-2700000')  # -45 deg en 1/60000e
        body_pr.set('vertOverflow', 'ellipsis')
        body_pr.set('wrap', 'square')
        body_pr.set('anchor', 'ctr')

        value_axis = chart.value_axis
        value_labels = value_axis.tick_labels
        value_labels.font.name = 'Arial'
        value_labels.font.size = Pt(9)

        legend = chart.legend
        if legend is not None:
            legend.font.name = 'Arial'
            legend.font.size = Pt(10)
            legend.include_in_layout = False

        print(f"  [OK] Slide '{slide_title}' créée avec {len(centres)} centre(s) et {len(metrics)} métrique(s)")
    
    else:
        # Pas de données pour cette organisation
        no_data_shape = slide.shapes.add_textbox(
            Inches(2), Inches(3), Inches(6), Inches(2)
        )
        no_data_frame = no_data_shape.text_frame
        no_data_frame.text = "Aucune donnée disponible pour cette organisation"
        
        no_data_paragraph = no_data_frame.paragraphs[0]
        no_data_paragraph.alignment = PP_ALIGN.CENTER
        no_data_font = no_data_paragraph.font
        no_data_font.name = 'Arial'
        no_data_font.size = Pt(20)
        no_data_font.color.rgb = RGBColor(190, 100, 100)
        
        print(f"  [INFO] Slide '{slide_title}' créée sans données")
    
    return slide

def main():
    """Créer PowerPoint final avec les 15 slides spécifiques"""
    
    print("=== Génération PowerPoint Final (15 slides spécifiques) ===")
    
    # Créer présentation
    prs = Presentation()
    
    # Slide titre
    title_slide_layout = prs.slide_layouts[0]
    title_slide = prs.slides.add_slide(title_slide_layout)
    title = title_slide.shapes.title
    subtitle = title_slide.placeholders[1]
    
    title.text = "Rapport CRM - Données Spécifiques"
    subtitle.text = "Slides générées avec données extraites du CRM"
    
    # Traitement de chaque slide configurée
    created_slides = 1  # Compteur incluant la slide titre
    for i, config in enumerate(SLIDE_CONFIG, 1):
        print(f"\nSlide {i+1}: {config['title']}")
        
        # Trouver le fichier SVG
        svg_file = find_svg_file(config['folder'], config['file_pattern'])
        if not svg_file:
            print(f"  [ERREUR] Fichier SVG non trouvé: {config['file_pattern']} dans {config['folder']}")
            continue
        
        print(f"  Fichier trouvé: {os.path.basename(svg_file)}")
        
        # Parser le SVG
        try:
            titre, legendes, valeurs, structure = extract_highcharts_data(svg_file)
            if not structure:
                print(f"  [WARNING] Aucune donnée extraite")
                # Créer slide vide
                slide_layout = prs.slide_layouts[5]
                slide = prs.slides.add_slide(slide_layout)
                title_shape = slide.shapes.title
                title_shape.text = config['title']
                created_slides += 1
                continue
                
        except Exception as e:
            print(f"  [ERREUR] Parsing échoué: {e}")
            # Créer slide vide
            slide_layout = prs.slide_layouts[5] 
            slide = prs.slides.add_slide(slide_layout)
            title_shape = slide.shapes.title
            title_shape.text = config['title']
            created_slides += 1
            continue
        
        # Créer la slide avec données
        slide = create_slide_with_specific_data(prs, config['title'], titre, structure, config['org'], i+1)
        if slide:
            created_slides += 1
    
    # Sauvegarder dans le répertoire parent (root)
    output_file = os.path.join(os.path.dirname(__file__), "..", "Rapport_CRM_Final_Donnees_Specifiques.pptx")
    prs.save(output_file)
    
    print(f"\n" + "="*70)
    print(f"[OK] PowerPoint créé: {output_file}")
    print(f"[INFO] Slides configurées: {len(SLIDE_CONFIG)}")
    print(f"[INFO] Total slides: {len(prs.slides)} (incluant titre)")
    print(f"[INFO] DONNÉES EXTRAITES DU CRM!")
    print("="*70)

    # Envoi automatique par email si configuré
    try:
        print("EMAIL: Tentative d'envoi automatique du rapport par email...")
        import sys
        sys.path.append(os.path.join(os.path.dirname(__file__), '..'))

        from email_sender.send_report import CRMReportSender
        sender = CRMReportSender()

        # Vérifier si la configuration email est valide
        from email_config.email_settings import validate_config
        config_errors = validate_config()

        if not config_errors:
            success = sender.send_report()
            if success:
                print("SUCCESS: Rapport envoyé par email avec succès!")
            else:
                print("ERROR: Échec de l'envoi par email")
        else:
            print("WARNING: Configuration email incomplète, envoi ignoré")
            print("   Pour activer l'envoi automatique, configurez email_config/email_settings.py")

    except ImportError:
        print("WARNING: Module d'envoi email non trouvé, génération seule effectuée")
    except Exception as e:
        print(f"WARNING: Erreur lors de l'envoi email: {e}")

    return output_file

if __name__ == "__main__":
    main()