#!/usr/bin/env python3
"""
Script pour générer le rapport PowerPoint PERSPECTIVIA
à partir des graphiques générés par t.py
"""

import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import glob

# Fix console encoding for Windows
if sys.platform == 'win32':
    sys.stdout.reconfigure(encoding='utf-8')

# =============================================================================
# CONFIGURATION DYNAMIQUE
# =============================================================================

# Utiliser des chemins relatifs au script
base_directory = os.path.dirname(os.path.abspath(__file__))
graph_directory = os.path.join(base_directory, 'output', 'Graphiques')
output_directory = base_directory
output_filename = 'Rapport_PERSPECTIVIA.pptx'

# =============================================================================
# CREATE POWERPOINT PRESENTATION
# =============================================================================

def create_presentation_with_graphs(graph_dir, output_dir, output_file):
    """
    Create a PowerPoint presentation with all graphs from the specified directory
    """
    try:
        # Create a presentation object
        prs = Presentation()
        prs.slide_width = Inches(10)  # Standard 16:9 ratio width
        prs.slide_height = Inches(7.5)  # Standard 16:9 ratio height

        # Add title slide
        title_slide_layout = prs.slide_layouts[0]  # Title slide layout
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]

        title.text = "Rapport d'Analyse PERSPECTIVIA"
        subtitle.text = "Visualisations des Formations par Vague\nGénéré automatiquement"

        # Get all PNG files from the graph directory
        all_graph_files = glob.glob(os.path.join(graph_dir, '*.png'))

        if not all_graph_files:
            print(f"⚠ WARNING: No PNG files found in {graph_dir}")
            return

        # Define graph order and titles (alphabetical order, Relances_par_Personne is last)
        graph_order = [
            'CA_par_Catégorie_Toutes_Vagues.png',
            'Paiements_par_Vague.png',
            'PROMO_Reel_par_Vague.png',
            'Statut_Formations_par_Vague.png',
            'Statuts_Intermédiaires_Potentiel.png',
            'Statuts_Intermédiaires_Previsionnel.png',
            'Statuts_Intermédiaires_Reel.png',
            'Relances_par_Personne.png'  # LAST SLIDE
        ]

        graph_info = {
            'Paiements_par_Vague.png': 'Paiements Réels par Vague et Type de Paiement',
            'Statut_Formations_par_Vague.png': 'Répartition des Formations par Vague et État',
            'CA_par_Catégorie_Toutes_Vagues.png': 'Chiffre d\'Affaires par Catégorie et par Vague',
            'Statuts_Intermédiaires_Reel.png': 'Statuts Intermédiaires - Réél',
            'Statuts_Intermédiaires_Previsionnel.png': 'Statuts Intermédiaires - Prévisionnel',
            'Statuts_Intermédiaires_Potentiel.png': 'Statuts Intermédiaires - Potentiel',
            'PROMO_Reel_par_Vague.png': 'Répartition des Formations Réelles par PROMO',
            'Relances_par_Personne.png': 'Suivi des Relances par Collaborateur'
        }

        # Sort graph files according to defined order
        graph_files = []
        for graph_name in graph_order:
            graph_path = os.path.join(graph_dir, graph_name)
            if os.path.exists(graph_path):
                graph_files.append(graph_path)
            else:
                print(f"⚠ WARNING: Expected graph not found: {graph_name}")

        print(f"\n📊 Found {len(graph_files)} graphs to add to presentation:")
        for graph_file in graph_files:
            print(f"   - {os.path.basename(graph_file)}")

        # Add a slide for each graph
        for graph_file in graph_files:
            # Use blank layout for full control
            blank_slide_layout = prs.slide_layouts[6]  # Blank layout
            slide = prs.slides.add_slide(blank_slide_layout)

            # Get graph filename
            graph_filename = os.path.basename(graph_file)

            # Get title from dictionary or use filename
            slide_title = graph_info.get(graph_filename, graph_filename.replace('.png', '').replace('_', ' '))

            # Add title text box at the top
            left = Inches(0.5)
            top = Inches(0.3)
            width = Inches(9)
            height = Inches(0.6)

            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.text = slide_title

            # Format title
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(24)
            p.font.bold = True
            p.font.name = 'Calibri'

            # Add the graph image
            # Calculate dimensions to fit well on slide
            img_left = Inches(0.5)
            img_top = Inches(1.2)
            img_width = Inches(9)

            # Réduire de 10% la largeur et remonter de 15% la position verticale pour "Répartition des Formations par Vague et État"
            if graph_filename == 'Statut_Formations_par_Vague.png':
                img_width = Inches(9 * 0.9)  # Réduction de 10% en largeur
                # Centrer le graphique réduit horizontalement
                img_left = Inches(0.5 + (9 * 0.1) / 2)
                # Remonter le graphique de 15% (diminuer img_top)
                img_top = Inches(1.2 * 0.85)  # Position verticale réduite de 15%

            # Réduire le graphique CA et le positionner à gauche pour laisser place aux annotations
            if graph_filename == 'CA_par_Catégorie_Toutes_Vagues.png':
                img_width = Inches(6.5)  # Graphique plus étroit
                img_left = Inches(0.3)   # Collé à gauche

            try:
                pic = slide.shapes.add_picture(graph_file, img_left, img_top, width=img_width)
                print(f"   ✓ Added slide: {slide_title}")

                # Add text annotations for CA graph
                if graph_filename == 'CA_par_Catégorie_Toutes_Vagues.png':
                    # Add three compact text boxes on the right side
                    text_left = Inches(7.0)
                    text_width = Inches(2.7)
                    text_height = Inches(0.8)

                    # Potentiel annotation (top)
                    txBox1 = slide.shapes.add_textbox(text_left, Inches(1.3), text_width, text_height)
                    tf1 = txBox1.text_frame
                    tf1.text = "Potentiel : lister les statuts=> Trésorerie incertain, prévoir 50% de pertes"
                    tf1.word_wrap = True
                    p1 = tf1.paragraphs[0]
                    p1.font.size = Pt(9)
                    p1.font.name = 'Calibri'
                    # Add border
                    txBox1.line.color.rgb = RGBColor(0, 0, 0)
                    txBox1.line.width = Pt(1)
                    txBox1.fill.solid()
                    txBox1.fill.fore_color.rgb = RGBColor(255, 255, 255)

                    # Prévisionnel annotation (middle)
                    txBox2 = slide.shapes.add_textbox(text_left, Inches(2.3), text_width, text_height)
                    tf2 = txBox2.text_frame
                    tf2.text = "Prévisionnel : attente de prise en charge=> Trésorerie prévisonnel a court terme"
                    tf2.word_wrap = True
                    p2 = tf2.paragraphs[0]
                    p2.font.size = Pt(9)
                    p2.font.name = 'Calibri'
                    # Add border
                    txBox2.line.color.rgb = RGBColor(0, 0, 0)
                    txBox2.line.width = Pt(1)
                    txBox2.fill.solid()
                    txBox2.fill.fore_color.rgb = RGBColor(255, 255, 255)

                    # Réél annotation (bottom)
                    txBox3 = slide.shapes.add_textbox(text_left, Inches(3.3), text_width, text_height)
                    tf3 = txBox3.text_frame
                    tf3.text = "Réél : lister les catégories=> Trésorerie certain à court terme"
                    tf3.word_wrap = True
                    p3 = tf3.paragraphs[0]
                    p3.font.size = Pt(9)
                    p3.font.name = 'Calibri'
                    # Add border
                    txBox3.line.color.rgb = RGBColor(0, 0, 0)
                    txBox3.line.width = Pt(1)
                    txBox3.fill.solid()
                    txBox3.fill.fore_color.rgb = RGBColor(255, 255, 255)

                    print(f"   ✓ Added text annotations for CA graph")

            except Exception as e:
                print(f"   ✗ ERROR adding image {graph_filename}: {e}")

        # Save the presentation
        output_path = os.path.join(output_dir, output_file)
        prs.save(output_path)

        print(f"\n✅ PowerPoint presentation created successfully!")
        print(f"📄 Output file: {output_path}")
        print(f"📊 Total slides: {len(prs.slides)} (1 title + {len(graph_files)} graphs)")

    except Exception as e:
        print(f"❌ ERROR creating PowerPoint presentation: {e}")
        import traceback
        traceback.print_exc()


# =============================================================================
# MAIN EXECUTION
# =============================================================================

if __name__ == "__main__":
    print("=" * 70)
    print("  POWERPOINT PRESENTATION GENERATOR")
    print("=" * 70)
    print(f"📂 Graph directory: {graph_directory}")
    print(f"📂 Output directory: {output_directory}")

    # Check if graph directory exists
    if not os.path.exists(graph_directory):
        print(f"\n❌ ERROR: Graph directory not found: {graph_directory}")
        print(f"Creating directory: {graph_directory}")
        os.makedirs(graph_directory, exist_ok=True)
    else:
        create_presentation_with_graphs(graph_directory, output_directory, output_filename)

    print("\n" + "=" * 70)
    print("✅ PROCESSING COMPLETE")
    print("=" * 70)
